from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
MATERIAL_ROOT = SCRIPT_DIR.parent
OUTPUT_ROOT = MATERIAL_ROOT / "outputs"
PPTX_PATH = OUTPUT_ROOT / "AI智能无人洗车仿真系统阶段性成果汇报.pptx"
SPEAKER_NOTES_PATH = OUTPUT_ROOT / "AI智能无人洗车仿真系统阶段性成果汇报_讲稿.md"

WIDE_WIDTH = Inches(13.333)
WIDE_HEIGHT = Inches(7.5)

COLORS = {
    "navy": RGBColor(20, 45, 78),
    "blue": RGBColor(37, 99, 235),
    "cyan": RGBColor(14, 165, 233),
    "slate": RGBColor(71, 85, 105),
    "light": RGBColor(241, 245, 249),
    "line": RGBColor(203, 213, 225),
    "white": RGBColor(255, 255, 255),
    "green": RGBColor(16, 133, 111),
    "orange": RGBColor(217, 119, 6),
    "red": RGBColor(190, 18, 60),
    "dark": RGBColor(15, 23, 42),
}


def set_run_font(run, size: int, color: RGBColor | None = None, bold: bool = False) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def set_text(shape, text: str, size: int, color: RGBColor, bold: bool = False) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.06)
    frame.margin_bottom = Inches(0.06)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    set_run_font(run, size, color, bold)


def add_textbox(slide, left, top, width, height, text, size=18, color=None, bold=False):
    color = color or COLORS["dark"]
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    set_run_font(run, size, color, bold)
    return shape


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.55), Inches(0.32), Inches(11.9), Inches(0.5), title, 26, COLORS["navy"], True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.56), Inches(0.92), Inches(1.4), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["blue"]
    accent.line.fill.background()
    if subtitle:
        add_textbox(slide, Inches(0.56), Inches(1.02), Inches(11.7), Inches(0.35), subtitle, 12, COLORS["slate"], False)


def add_footer(slide, index: int) -> None:
    add_textbox(slide, Inches(0.58), Inches(7.05), Inches(5.4), Inches(0.22), "AI Car Wash Simulation Demo / Stage 1-3 Result", 8, COLORS["slate"])
    add_textbox(slide, Inches(12.05), Inches(7.05), Inches(0.75), Inches(0.22), f"{index:02d}/12", 8, COLORS["slate"])


def add_card(slide, left, top, width, height, title: str, body: str = "", color=None) -> None:
    color = color or COLORS["blue"]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.14)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.1)
    frame.margin_bottom = Inches(0.08)
    title_para = frame.paragraphs[0]
    run = title_para.add_run()
    run.text = title
    set_run_font(run, 14, color, True)
    if body:
        body_para = frame.add_paragraph()
        body_para.space_before = Pt(4)
        body_run = body_para.add_run()
        body_run.text = body
        set_run_font(body_run, 10, COLORS["slate"], False)


def add_bullets(slide, left, top, width, height, bullets: list[str], size=15) -> None:
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.level = 0
        paragraph.space_after = Pt(5)
        run = paragraph.add_run()
        run.text = f"• {item}"
        set_run_font(run, size, COLORS["dark"], False)


def add_flow(slide, labels: list[str], left, top, width, height, columns: int) -> None:
    gap = Inches(0.12)
    box_w = (width - gap * (columns - 1)) / columns
    box_h = height
    for i, label in enumerate(labels):
        x = left + (box_w + gap) * i
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["light"] if i % 2 == 0 else RGBColor(232, 240, 255)
        shape.line.color.rgb = COLORS["line"]
        set_text(shape, label, 12, COLORS["navy"], True)
        if i < len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w - Inches(0.03), top + box_h / 2 - Inches(0.1), Inches(0.28), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["blue"]
            arrow.line.fill.background()


def add_metric(slide, left, top, width, height, label: str, value: str, color=None) -> None:
    color = color or COLORS["blue"]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.1)
    frame.margin_right = Inches(0.1)
    frame.margin_top = Inches(0.08)
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = value
    set_run_font(run, 22, color, True)
    label_para = frame.add_paragraph()
    label_para.alignment = PP_ALIGN.CENTER
    label_run = label_para.add_run()
    label_run.text = label
    set_run_font(label_run, 9, COLORS["slate"], False)


def add_section_label(slide, left, top, text: str, color=None) -> None:
    color = color or COLORS["green"]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.35), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    set_text(shape, text, 9, COLORS["white"], True)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = WIDE_WIDTH
    prs.slide_height = WIDE_HEIGHT

    # 1
    slide = blank_slide(prs)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["navy"]
    add_textbox(slide, Inches(0.75), Inches(1.3), Inches(11.6), Inches(0.9), "AI智能无人洗车仿真系统阶段性成果汇报", 34, COLORS["white"], True)
    add_textbox(slide, Inches(0.78), Inches(2.35), Inches(10.8), Inches(0.45), "从车辆识别到洗车仿真、可视化演示和客户材料的阶段性软件闭环", 18, RGBColor(214, 226, 245))
    add_flow(slide, ["车辆识别", "洗车仿真", "覆盖评估", "可视化展示"], Inches(0.82), Inches(3.45), Inches(8.9), Inches(0.74), 4)
    add_textbox(slide, Inches(0.82), Inches(6.62), Inches(5.4), Inches(0.28), "AI Car Wash Simulation Demo / Stage 1-3 Result", 11, RGBColor(214, 226, 245))

    # 2
    slide = blank_slide(prs)
    add_title(slide, "行业痛点", "传统无人洗车方案如果直接进入硬件联调，试错成本高、流程不可解释、方案沟通困难。")
    cards = [
        ("喷嘴布局", "清洗区域和覆盖关系难提前验证。"),
        ("机械/PLC联调", "真实设备试错成本高，周期长。"),
        ("客户沟通", "系统如何工作不容易被直观看懂。"),
        ("软件闭环", "缺少可演示、可复盘的仿真原型。"),
    ]
    for i, (title, body) in enumerate(cards):
        add_card(slide, Inches(0.75 + i * 3.05), Inches(2.15), Inches(2.65), Inches(2.45), title, body, COLORS["blue"])
    add_bullets(slide, Inches(1.0), Inches(5.25), Inches(10.9), Inches(0.9), ["核心问题不是单点技术，而是识别、策略、路径、控制和展示缺少可解释的前置验证。"], 17)
    add_footer(slide, 2)

    # 3
    slide = blank_slide(prs)
    add_title(slide, "先软件仿真，再硬件联调", "先把识别、策略、空间、喷嘴、流程、路径、覆盖率和展示做成软件闭环。")
    add_flow(slide, ["输入车辆图片", "AI识别", "洗车仿真", "覆盖评估", "可视化展示", "后续硬件联调"], Inches(0.65), Inches(2.15), Inches(12.0), Inches(0.85), 6)
    add_card(slide, Inches(1.0), Inches(4.0), Inches(3.2), Inches(1.4), "先验证逻辑", "让流程和策略先在软件中跑通。", COLORS["green"])
    add_card(slide, Inches(5.0), Inches(4.0), Inches(3.2), Inches(1.4), "再接硬件", "后续进入真实路径规划和 PLC 联调。", COLORS["blue"])
    add_card(slide, Inches(9.0), Inches(4.0), Inches(3.2), Inches(1.4), "降低风险", "减少直接上设备的盲目试错。", COLORS["orange"])
    add_footer(slide, 3)

    # 4
    slide = blank_slide(prs)
    add_title(slide, "系统总体架构", "车辆识别结果以 JSON 形式驱动后续洗车策略和仿真模块。")
    flow = [
        "车辆图片",
        "vehicle_type_result.json",
        "wash_strategy_plan.json",
        "space_model_report.json",
        "nozzle_coverage_plan.json",
        "wash_flow_run.json",
        "abstract_nozzle_path_plan.json",
        "coverage_report.json",
        "2D / Timeline / Customer Showcase",
    ]
    for i, label in enumerate(flow):
        row = i // 3
        col = i % 3
        add_card(slide, Inches(0.75 + col * 4.05), Inches(1.55 + row * 1.45), Inches(3.55), Inches(0.95), label, "", COLORS["blue"])
    add_bullets(slide, Inches(0.95), Inches(6.25), Inches(11.0), Inches(0.55), ["模块边界清楚，后续可以逐步替换成更真实的算法、参数或硬件接口。"], 14)
    add_footer(slide, 4)

    # 5
    slide = blank_slide(prs)
    add_title(slide, "阶段1：车辆识别", "这是整个洗车策略选择的入口。")
    items = [("输入车辆图片", "test_car.jpg / demo input"), ("检测车辆", "YOLO bbox"), ("裁切车辆区域", "crop"), ("分类车型", "sedan / suv / mpv"), ("输出JSON", "vehicle_type_result.json")]
    for i, (title, body) in enumerate(items):
        add_card(slide, Inches(0.65 + i * 2.5), Inches(2.0), Inches(2.15), Inches(1.45), title, body, COLORS["green"])
    add_bullets(slide, Inches(1.0), Inches(4.6), Inches(10.9), Inches(1.2), [
        "AI识别结果不是停留在页面展示，而是继续传给 aicar_sim。",
        "当前三分类模型是小样本 demo，不代表商业级识别精度。"
    ], 16)
    add_footer(slide, 5)

    # 6
    slide = blank_slide(prs)
    add_title(slide, "阶段2：洗车仿真主链路", "把车型结果转成策略、空间、喷嘴、流程、路径和覆盖率报告。")
    sim_items = ["车型匹配", "洗车策略", "车辆包络", "洗车房空间", "喷嘴模型", "流程状态机", "抽象路径", "覆盖率报告"]
    for i, label in enumerate(sim_items):
        add_card(slide, Inches(0.65 + (i % 4) * 3.05), Inches(1.55 + (i // 4) * 1.35), Inches(2.55), Inches(0.9), label, "", COLORS["blue"])
    metrics = [("estimated_total_seconds", "141s"), ("segment_count", "22"), ("point_count", "112"), ("coverage", "92%")]
    for i, (label, val) in enumerate(metrics):
        add_metric(slide, Inches(1.0 + i * 3.0), Inches(5.1), Inches(2.35), Inches(0.9), label, val, COLORS["green"])
    add_footer(slide, 6)

    # 7
    slide = blank_slide(prs)
    add_title(slide, "阶段3：可视化与动画演示", "让技术结果变成客户能看懂的演示。")
    cards = [
        ("2D俯视图", "洗车房、车辆包络、路径点。"),
        ("侧视图", "高度与空间关系。"),
        ("时间轴动画", "流程状态和当前区域高亮。"),
        ("客户展示页", "项目价值、边界和路线。"),
    ]
    for i, (title, body) in enumerate(cards):
        add_card(slide, Inches(0.85 + i * 3.0), Inches(2.0), Inches(2.55), Inches(2.0), title, body, COLORS["cyan"])
    add_bullets(slide, Inches(1.0), Inches(5.15), Inches(10.8), Inches(0.9), ["展示层不改变底层仿真逻辑，它的作用是让结果可看、可讲、可复盘。"], 16)
    add_footer(slide, 7)

    # 8
    slide = blank_slide(prs)
    add_title(slide, "当前 Demo 指标", "当前样例已完成识别、策略、路径和覆盖率估算闭环。")
    metrics = [
        ("vehicle_type", "sedan"),
        ("wash_profile", "standard_sedan"),
        ("estimated_total_seconds", "141s"),
        ("state_count", "10"),
        ("segment_count", "22"),
        ("point_count", "112"),
        ("coverage", "92%"),
        ("coverage_pass", "true"),
    ]
    for i, (label, val) in enumerate(metrics):
        add_metric(slide, Inches(0.75 + (i % 4) * 3.05), Inches(1.75 + (i // 4) * 1.65), Inches(2.55), Inches(1.05), label, val, COLORS["blue"] if i < 4 else COLORS["green"])
    add_footer(slide, 8)

    # 9
    slide = blank_slide(prs)
    add_title(slide, "客户价值", "从概念图升级为可运行软件原型。")
    value_cards = [
        ("对设备厂商", "降低试错成本\n提前验证喷嘴布局和流程逻辑"),
        ("对运营方", "更容易理解无人洗车流程\n后续可沉淀运行数据"),
        ("对研发团队", "软件闭环清晰\n后续可逐步接硬件"),
        ("对客户演示", "从概念图升级为可运行软件原型"),
    ]
    for i, (title, body) in enumerate(value_cards):
        add_card(slide, Inches(0.85 + (i % 2) * 6.05), Inches(1.8 + (i // 2) * 1.85), Inches(5.35), Inches(1.35), title, body, COLORS["green"])
    add_footer(slide, 9)

    # 10
    slide = blank_slide(prs)
    add_title(slide, "技术边界", "当前是软件仿真闭环，不是商用硬件控制系统。")
    add_section_label(slide, Inches(0.9), Inches(1.5), "当前可以说", COLORS["green"])
    add_bullets(slide, Inches(0.9), Inches(1.95), Inches(5.3), Inches(3.4), [
        "已完成软件仿真闭环",
        "已完成阶段性可视化 Demo",
        "已完成抽象路径和覆盖率报告",
        "可作为后续路径规划和硬件联调基础",
    ], 14)
    add_section_label(slide, Inches(7.0), Inches(1.5), "当前不能说", COLORS["red"])
    add_bullets(slide, Inches(7.0), Inches(1.95), Inches(5.3), Inches(3.4), [
        "已经实现真实无人洗车控制",
        "已经验证真实清洗效果",
        "已经完成PLC联调",
        "已经完成商业级AI识别",
        "已经完成真实机械轨迹规划",
    ], 14)
    add_footer(slide, 10)

    # 11
    slide = blank_slide(prs)
    add_title(slide, "后续路线：阶段4/5/6", "从软件闭环走向真实系统。")
    roadmap = [
        ("阶段4", "真实路径规划与运动约束", "速度、加速度、边界、碰撞、安全约束"),
        ("阶段5", "PLC/硬件联调", "I/O、状态机、急停、喷嘴/泵/风机控制"),
        ("阶段6", "商业化后台与设备管理", "设备管理、订单用户、洗车记录、远程运维"),
    ]
    for i, (stage, title, body) in enumerate(roadmap):
        add_metric(slide, Inches(1.0 + i * 4.0), Inches(1.7), Inches(1.2), Inches(0.8), stage, str(i + 4), COLORS["blue"])
        add_card(slide, Inches(0.8 + i * 4.0), Inches(2.75), Inches(3.25), Inches(2.1), title, body, COLORS["blue"])
    add_footer(slide, 11)

    # 12
    slide = blank_slide(prs)
    add_title(slide, "总结与合作方向", "当前项目已经具备可运行、可展示、可讲解的软件原型基础。")
    add_card(slide, Inches(0.85), Inches(1.65), Inches(5.4), Inches(3.6), "阶段性成果", "完成从识别、仿真、路径、覆盖率到可视化展示的阶段性闭环。", COLORS["green"])
    add_bullets(slide, Inches(7.0), Inches(1.75), Inches(5.2), Inches(3.2), [
        "继续打磨客户演示",
        "引入真实车辆/洗车机参数",
        "推进真实路径规划",
        "推进PLC/硬件联调",
        "探索商业化系统",
    ], 15)
    add_footer(slide, 12)

    return prs


SPEAKER_NOTES = [
    ("页1", "封面", "大家好，这份汇报介绍的是 AI 智能无人洗车仿真系统的阶段性成果。当前我们已经完成从车辆识别到洗车仿真、可视化展示和客户材料的一个软件闭环。需要先说明，这不是最终商用设备发布，而是一个可运行、可展示、可继续深化的软件原型。", "约45秒"),
    ("页2", "行业痛点", "无人洗车如果一开始直接进入硬件联调，问题会集中在后期暴露，包括喷嘴布局、洗车流程、机械运动、PLC控制和客户沟通。我们希望先用软件方式把这些逻辑讲清楚、跑起来，减少后续硬件试错。", "约50秒"),
    ("页3", "项目思路", "我们的思路是先软件仿真，再硬件联调。先输入车辆图片，通过AI识别车型，再生成洗车仿真、覆盖评估和可视化展示。等这个闭环稳定后，再进入真实路径规划和PLC硬件联调。", "约55秒"),
    ("页4", "系统总体架构", "系统通过 JSON 把各个模块串起来。车辆识别输出 vehicle_type_result.json，后续依次生成洗车策略、空间模型、喷嘴覆盖、流程状态机、抽象路径和覆盖率报告，最后进入2D、时间轴和客户展示页面。", "约60秒"),
    ("页5", "阶段1：车辆识别", "阶段1解决的是车辆识别入口问题。系统能从车辆图片中检测车辆、裁切车辆区域，并分类为 sedan、suv 或 mpv，然后输出标准 JSON。这个结果会继续交给 aicar_sim 使用，作为洗车策略选择入口。", "约55秒"),
    ("页6", "阶段2：洗车仿真主链路", "阶段2把车型结果变成可检查的仿真输出，包括车型匹配、洗车策略、车辆包络、洗车房空间、喷嘴模型、流程状态机、抽象路径和覆盖率报告。当前样例总时长141秒，路径段22个，路径点112个，覆盖估算为92%。", "约65秒"),
    ("页7", "阶段3：可视化与动画演示", "阶段3把阶段2的 JSON 结果变成客户能看懂的页面，包括2D俯视图、侧视图、时间轴动画和客户展示页。它不改变底层逻辑，但让结果更容易汇报、讲解和复盘。", "约50秒"),
    ("页8", "当前 Demo 指标", "这里是当前样例的关键指标：车型是 sedan，策略是 standard_sedan，总流程141秒，状态数10个，路径段22个，路径点112个，覆盖率估算92%，并且 coverage_pass 为 true。注意这里的覆盖率是软件估算，不是实际洗净率。", "约55秒"),
    ("页9", "客户价值", "这个原型的价值在于把概念图变成可运行软件。对设备厂商，它能降低试错成本；对运营方，它让流程更容易理解；对研发团队，它让模块边界清晰；对客户演示，它能直观看到识别、策略、路径、覆盖和动画。", "约60秒"),
    ("页10", "技术边界", "这里要特别强调边界。当前可以说已经完成软件仿真闭环和阶段性可视化Demo，可作为后续路径规划和硬件联调基础。但不能说已经实现真实无人洗车控制，也不能说已经验证真实清洗效果或完成PLC联调。", "约60秒"),
    ("页11", "后续路线", "后续建议分三步走。阶段4做真实路径规划和运动约束，把抽象路径变得更接近机械运动。阶段5进入PLC和硬件联调。阶段6再扩展商业化后台，包括设备管理、订单用户、洗车记录和远程运维。", "约60秒"),
    ("页12", "总结与合作方向", "总结来说，当前项目已经完成从识别、仿真、路径、覆盖率到可视化展示的阶段性闭环。后续可以继续打磨客户演示，也可以引入真实车辆和洗车机参数，推进真实路径规划、PLC硬件联调和商业化系统。", "约55秒"),
]


def build_speaker_notes() -> str:
    lines = [
        "# AI智能无人洗车仿真系统阶段性成果汇报_讲稿",
        "",
        "总讲解时间建议控制在 8-12 分钟。",
        "",
    ]
    for page, title, script, duration in SPEAKER_NOTES:
        lines.extend(
            [
                f"## {page}：{title}",
                "",
                f"- 页面标题：{title}",
                f"- 预计讲解时间：{duration}",
                "- 讲解话术：",
                "",
                script,
                "",
            ]
        )
    return "\n".join(lines)


def main() -> None:
    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(PPTX_PATH)
    SPEAKER_NOTES_PATH.write_text(build_speaker_notes(), encoding="utf-8")
    print(f"PPTX saved: {PPTX_PATH}")
    print(f"speaker notes saved: {SPEAKER_NOTES_PATH}")
    print(f"slide_count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
